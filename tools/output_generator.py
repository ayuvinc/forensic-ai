"""Output generator — .docx and .pptx generation from markdown and storyboard schemas."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Optional


class OutputGenerator:
    """Generate .docx and .pptx deliverables from content."""

    def generate_docx(
        self,
        content_md: str,
        output_path: str | Path,
        template_path: Optional[str | Path] = None,
        workflow_type: Optional[str] = None,
        case_id: Optional[str] = None,
        firm_name: str = "GoodWork Forensic Consulting",
    ) -> Path:
        """Generate .docx from markdown content.

        Resolution order (TPL-02):
          1. explicit template_path (if provided and exists)
          2. TemplateManager.resolve(workflow_type) when workflow_type is set
          3. plain python-docx Document() — no template

        Applies GW_ named styles from template when available. Falls back
        silently to basic styling if a named style is absent.
        Logs a template_resolved audit event to cases/{case_id}/audit_log.jsonl
        when case_id is provided.
        """
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)

        resolved_template: Optional[Path] = None
        fallback = True

        # Step 1: explicit path wins
        if template_path and Path(template_path).exists():
            resolved_template = Path(template_path)
            fallback = False
        elif workflow_type:
            # Step 2: resolve via TemplateManager
            try:
                from tools.template_manager import TemplateManager
                resolved_template = TemplateManager().resolve(workflow_type)
                fallback = False
            except Exception:
                resolved_template = None

        if resolved_template:
            doc = Document(str(resolved_template))
        else:
            # Step 3: plain document
            doc = Document()
            style = doc.styles["Normal"]
            style.font.name = "Calibri"
            style.font.size = Pt(11)

        # Detect which GW_ named styles are available in the loaded template
        available_styles = {s.name for s in doc.styles} if resolved_template else set()

        _md_to_docx(doc, content_md, available_styles)

        # Audit log — best-effort, non-blocking
        if case_id:
            try:
                from tools.file_tools import append_audit_event
                template_name = resolved_template.name if resolved_template else "none"
                append_audit_event(case_id, {
                    "event": "template_resolved",
                    "template": template_name,
                    "workflow_type": workflow_type or "",
                    "fallback": fallback,
                })
            except Exception:
                pass

        tmp = out.with_suffix(".tmp.docx")
        doc.save(tmp)
        import os
        os.replace(tmp, out)
        return out

    def generate_pptx(
        self,
        storyboard,   # DeckStoryboard
        output_path: str | Path,
        template_path: Optional[str | Path] = None,
        firm_name: str = "GoodWork Forensic Consulting",
    ) -> Path:
        """Generate .pptx from DeckStoryboard — one slide per SlideSpec."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)

        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
        else:
            prs = Presentation()

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        _set_text(slide.placeholders, 0, storyboard.deck_objective)
        _set_text(slide.placeholders, 1, f"{firm_name}\n{storyboard.audience}")

        # Content slides
        content_layout = prs.slide_layouts[1]
        for spec in storyboard.slides:
            slide = prs.slides.add_slide(content_layout)
            _set_text(slide.placeholders, 0, f"{spec.slide_number}. {spec.title}")

            bullets = []
            bullets.append(f"Key message: {spec.key_message}")
            bullets.extend(spec.content_bullets)
            if spec.evidence_needed:
                bullets.append(f"Evidence: {'; '.join(spec.evidence_needed[:2])}")

            _set_text(slide.placeholders, 1, "\n".join(bullets))

            # Speaker notes
            if spec.speaker_notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = spec.speaker_notes

        tmp = out.with_suffix(".tmp.pptx")
        prs.save(tmp)
        import os
        os.replace(tmp, out)
        return out


def _md_to_docx(doc, content: str, available_styles: set | None = None) -> None:
    """Convert markdown content to docx paragraphs.

    Applies GW_ named styles (TPL-02) when the loaded template provides them.
    Falls back to built-in heading/paragraph styles when GW_ styles are absent.
    available_styles: set of style names present in the document — None means unknown.
    """
    if available_styles is None:
        available_styles = set()

    def _has(style_name: str) -> bool:
        return style_name in available_styles

    def _add_para(text: str, gw_style: str, fallback_style: str) -> None:
        style = gw_style if _has(gw_style) else fallback_style
        try:
            doc.add_paragraph(text, style=style)
        except Exception:
            doc.add_paragraph(text)

    for line in content.split("\n"):
        line = line.rstrip()
        if line.startswith("# "):
            text = line[2:]
            if _has("GW_Title"):
                try:
                    doc.add_paragraph(text, style="GW_Title")
                except Exception:
                    doc.add_heading(text, level=1)
            else:
                doc.add_heading(text, level=1)
        elif line.startswith("## "):
            text = line[3:]
            if _has("GW_Heading1"):
                try:
                    doc.add_paragraph(text, style="GW_Heading1")
                except Exception:
                    doc.add_heading(text, level=2)
            else:
                doc.add_heading(text, level=2)
        elif line.startswith("### "):
            text = line[4:]
            if _has("GW_Heading2"):
                try:
                    doc.add_paragraph(text, style="GW_Heading2")
                except Exception:
                    doc.add_heading(text, level=3)
            else:
                doc.add_heading(text, level=3)
        elif line.startswith("- ") or line.startswith("* "):
            try:
                doc.add_paragraph(line[2:], style="List Bullet")
            except Exception:
                doc.add_paragraph(line[2:])
        elif line.startswith("| "):
            doc.add_paragraph(line)
        elif line.startswith("---"):
            doc.add_paragraph("_" * 50)
        elif line.strip():
            clean = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', line)
            clean = re.sub(r'_{1,2}([^_]+)_{1,2}', r'\1', clean)
            if _has("GW_Body"):
                try:
                    doc.add_paragraph(clean, style="GW_Body")
                except Exception:
                    doc.add_paragraph(clean)
            else:
                doc.add_paragraph(clean)
        else:
            doc.add_paragraph("")


def _set_text(placeholders, idx: int, text: str) -> None:
    """Safely set text on a pptx placeholder."""
    try:
        placeholders[idx].text = text
    except (KeyError, IndexError):
        pass
